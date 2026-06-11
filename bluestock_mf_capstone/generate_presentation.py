import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paths
BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = BASE_DIR / "Bluestock_MF_Presentation.pptx"

# Color Palette (Dark Theme matching Bluestock Dashboard)
BG_COLOR = RGBColor(10, 14, 39)       # #0A0E27 Navy
TEXT_LIGHT = RGBColor(255, 255, 255) # White
TEXT_MUTED = RGBColor(136, 153, 170) # Grey #8899AA
ACCENT_BLUE = RGBColor(30, 144, 255) # Blue #1E90FF
ACCENT_CYAN = RGBColor(0, 212, 255)  # Cyan #00D4FF
ACCENT_GREEN = RGBColor(0, 230, 118) # Green #00E676

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    return title_box

def main():
    print("Generating Bluestock MF Presentation PPTX (12 slides)...")
    prs = Presentation()
    # Change slide dimensions to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_slide_layout = prs.slide_layouts[6] # Blank layout
    
    # ------------------ SLIDE 1: TITLE SLIDE ------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide1)
    
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BLUESTOCK MUTUAL FUND ANALYTICS"
    p.font.name = "Arial"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "End-to-End ETL, Analytics, and Data Visualization Pipeline"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "Prepared by: Divvyansh Kudesiaa\nDate: June 2026\nCapstone Project Deliverable"
    p3.font.name = "Arial"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(30)
    
    # ------------------ SLIDE 2: PROBLEM & OBJECTIVE ------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide2)
    add_title(slide2, "Problem Statement & Objectives")
    
    tb = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Challenges in Mutual Fund Distribution Platforms:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(8)
    
    p_b1 = tf.add_paragraph()
    p_b1.text = "• Heterogeneous Data Sources: Hard to merge static PDFs/CSVs with live transactional streams."
    p_b1.font.size = Pt(14)
    p_b1.font.color.rgb = TEXT_LIGHT
    p_b1.space_after = Pt(6)
    
    p_b2 = tf.add_paragraph()
    p_b2.text = "• Risk Calculations: Lack of dynamic computations for Alpha, Beta, Sharpe, and Sortino metrics."
    p_b2.font.size = Pt(14)
    p_b2.font.color.rgb = TEXT_LIGHT
    p_b2.space_after = Pt(6)

    p_b3 = tf.add_paragraph()
    p_b3.text = "• Investor Personalization: Distribution models do not align with client risk appetites."
    p_b3.font.size = Pt(14)
    p_b3.font.color.rgb = TEXT_LIGHT
    p_b3.space_after = Pt(20)

    p_obj = tf.add_paragraph()
    p_obj.text = "Core Capstone Goals:"
    p_obj.font.bold = True
    p_obj.font.size = Pt(18)
    p_obj.font.color.rgb = ACCENT_GREEN
    p_obj.space_after = Pt(8)
    
    p_o1 = tf.add_paragraph()
    p_o1.text = "• Automate ETL pipeline, load structured relational star schema database, and validate integrity."
    p_o1.font.size = Pt(14)
    p_o1.font.color.rgb = TEXT_LIGHT
    
    # ------------------ SLIDE 3: DATA SOURCES ------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide3)
    add_title(slide3, "heterogeneous Data Sources")
    
    tb = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "The pipeline integrates 10 distinct datasets from Google Drive & live APIs:"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(12)
    
    items = [
        "1. Fund Master (CSV): AMC details, launch date, benchmark and expense ratios.",
        "2. NAV History (CSV): 46,000+ daily NAV prices across 40 schemes.",
        "3. Live NAV (API): Live daily values scraped from API.mfapi.in and portal.amfiindia.com.",
        "4. Industry AUM (CSV): Asset sizes and scheme counts for 90 AMCs.",
        "5. Category Inflows (CSV): 144 months of net flows across categories.",
        "6. Retail Transactions (CSV): 32,778 transaction records detailing investor age, state, and payment mode.",
        "7. Portfolio Holdings (CSV): Stock weights for each scheme."
    ]
    for item in items:
        p_item = tf.add_paragraph()
        p_item.text = f"• {item}"
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_MUTED
        p_item.space_after = Pt(6)

    # ------------------ SLIDE 4: ARCHITECTURE ------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide4)
    add_title(slide4, "Pipeline Architecture & ETL design")
    
    tb = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Extract-Transform-Load Pipeline Orchestration:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)
    
    steps = [
        ("1. Ingestion", "Downloads raw datasets from Google Drive and updates live NAVs from APIs."),
        ("2. Transformation", "Handles duplicates, standardises transaction types/KYC states, clips expense ratio outliers (0.1%-2.5%)."),
        ("3. Re-indexing", "Generates full calendar ranges per scheme and forward-fills missing weekend/holiday NAVs."),
        ("4. Star Schema", "Populates dim_fund, dim_date, fact_nav, fact_transactions, fact_performance, and fact_aum inside SQLite."),
        ("5. Visualization", "generate_dashboard.py constructs 4 PNG reports and combines them into Dashboard.pdf.")
    ]
    
    for step, desc in steps:
        p_step = tf.add_paragraph()
        p_step.text = f"• {step}: {desc}"
        p_step.font.size = Pt(13)
        p_step.font.color.rgb = TEXT_LIGHT
        p_step.space_after = Pt(6)

    # ------------------ SLIDE 5: EDA HIGHLIGHTS 1 ------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide5)
    add_title(slide5, "EDA Highlights: industry AUM & SIP trends")
    
    tb = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Industry AUM Expansion:"
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p_desc = tf.add_paragraph()
    p_desc.text = "• Market assets under management for the top 10 AMCs rose from ₹3.2L Cr to over ₹5.6L Cr.\n• SBI, ICICI, and HDFC represent the largest market shares, accounting for over 56% of total sector AUM."
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = TEXT_LIGHT
    
    tb2 = slide5.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "SIP Inflow Trends:"
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.size = Pt(18)
    p2.space_after = Pt(8)
    
    p2_desc = tf2.add_paragraph()
    p2_desc.text = "• Monthly SIP inflows grew from ~₹11,000 Crore to over ₹16,000 Crore.\n• Showcases stable recurring retail investments that resist direct market volatility."
    p2_desc.font.size = Pt(13)
    p2_desc.font.color.rgb = TEXT_LIGHT

    # ------------------ SLIDE 6: EDA HIGHLIGHTS 2 ------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide6)
    add_title(slide6, "EDA Highlights: Retail Demographics")
    
    tb = slide6.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key Insights from 32,000+ Retail Transactions:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)
    
    demographics = [
        "• Geographic Concentration: Maharashtra, Gujarat, Karnataka, Delhi, and Tamil Nadu comprise 60%+ of investment values.",
        "• Age Breakdown: The 26-35 age bracket forms the highest transaction count (tech-savvy young professionals).",
        "• Ticket Sizes: Younger investors prefer smaller recurring SIPs; older brackets (46-55, 56+) commit large lumpsum ticket sizes.",
        "• Payment Modes: UPI and Mandate transfers dominate SIP setups, while Net Banking remains preferred for lumpsums.",
        "• KYC Compliance: Overall KYC verification rate is at 92.4%, showing high compliance but leaving room for platform onboarding improvements."
    ]
    for dem in demographics:
        p_dem = tf.add_paragraph()
        p_dem.text = dem
        p_dem.font.size = Pt(13)
        p_dem.font.color.rgb = TEXT_LIGHT
        p_dem.space_after = Pt(8)

    # ------------------ SLIDE 7: PERFORMANCE METRICS 1 ------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide7)
    add_title(slide7, "Performance Metrics: Risk vs Return")
    
    # Left description
    tb = slide7.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Risk-Return Profiles:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(8)
    
    p_b1 = tf.add_paragraph()
    p_b1.text = "• Scatter mapping returns against standard deviation reveals equity categories reside in high-volatility zones."
    p_b1.font.size = Pt(13)
    p_b1.font.color.rgb = TEXT_LIGHT
    p_b1.space_after = Pt(6)
    
    p_b2 = tf.add_paragraph()
    p_b2.text = "• Benchmark comparison chart shows Axis Bluechip and SBI Bluechip outperforming Nifty 50 on a risk-adjusted basis (lower drawdowns)."
    p_b2.font.size = Pt(13)
    p_b2.font.color.rgb = TEXT_LIGHT
    
    # Right image
    chart_path = BASE_DIR / "benchmark_comparison.png"
    if chart_path.exists():
        slide7.shapes.add_picture(str(chart_path), Inches(5.0), Inches(1.3), Inches(4.5), Inches(3.5))

    # ------------------ SLIDE 8: PERFORMANCE METRICS 2 ------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide8)
    add_title(slide8, "Performance Metrics: Risk-Adjusted Ratios")
    
    # Left description
    tb = slide8.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Statistical Variables Computed:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(8)
    
    metrics = [
        "• Sharpe Ratio: Excess returns per unit of total risk.",
        "• Sortino Ratio: Volatility index based purely on negative downside returns.",
        "• Alpha: Active outperformance against benchmark indices.",
        "• Beta: Systemic risk sensitivity relative to the market index."
    ]
    for metric in metrics:
        p_met = tf.add_paragraph()
        p_met.text = metric
        p_met.font.size = Pt(12)
        p_met.font.color.rgb = TEXT_LIGHT
        p_met.space_after = Pt(6)
        
    # Right image
    sharpe_chart = BASE_DIR / "rolling_sharpe_chart.png"
    if sharpe_chart.exists():
        slide8.shapes.add_picture(str(sharpe_chart), Inches(5.0), Inches(1.3), Inches(4.5), Inches(3.5))

    # ------------------ SLIDE 9: DASHBOARD SCREENSHOTS 1 ------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide9)
    add_title(slide9, "Dashboard: Industry & Fund Overview")
    
    # Page 1 Image
    img1 = BASE_DIR / "dashboard" / "Page1_Industry_Overview.png"
    if img1.exists():
        slide9.shapes.add_picture(str(img1), Inches(0.5), Inches(1.3), Inches(4.2), Inches(2.6))
        tb = slide9.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(4.2), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Page 1: Industry Overview\nTotal AUM, SIP growth trends, and active scheme distributions."
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        
    # Page 2 Image
    img2 = BASE_DIR / "dashboard" / "Page2_Fund_Performance.png"
    if img2.exists():
        slide9.shapes.add_picture(str(img2), Inches(5.3), Inches(1.3), Inches(4.2), Inches(2.6))
        tb2 = slide9.shapes.add_textbox(Inches(5.3), Inches(4.0), Inches(4.2), Inches(1.2))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "Page 2: Fund Performance\nRisk scatters, top schemes scorecard, and NAV benchmark tracking."
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_LIGHT

    # ------------------ SLIDE 10: DASHBOARD SCREENSHOTS 2 ------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide10)
    add_title(slide10, "Dashboard: Investor & SIP Trends")
    
    # Page 3 Image
    img3 = BASE_DIR / "dashboard" / "Page3_Investor_Analytics.png"
    if img3.exists():
        slide10.shapes.add_picture(str(img3), Inches(0.5), Inches(1.3), Inches(4.2), Inches(2.6))
        tb = slide10.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(4.2), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Page 3: Investor Analytics\nState transaction volumes, product splits, and age preferences."
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        
    # Page 4 Image
    img4 = BASE_DIR / "dashboard" / "Page4_SIP_Market_Trends.png"
    if img4.exists():
        slide10.shapes.add_picture(str(img4), Inches(5.3), Inches(1.3), Inches(4.2), Inches(2.6))
        tb2 = slide10.shapes.add_textbox(Inches(5.3), Inches(4.0), Inches(4.2), Inches(1.2))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = "Page 4: SIP & Market Trends\nSIP correlation with index values and monthly category net flows."
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_LIGHT

    # ------------------ SLIDE 11: KEY FINDINGS & RECOMMENDATIONS ------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide11)
    add_title(slide11, "Strategic Recommendations")
    
    tb = slide11.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Actionable Insights for Mutual Fund Platforms:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(10)
    
    recs = [
        "1. Focus on Tier-2/3 Geographic Expansion: Build UPI Auto-Pay integrations to ease digital boarding.",
        "2. Automated Risk profiling: Use CLI recommender models to automate fund suggestions by Sharpe ratio.",
        "3. Highlight Expense Ratios: Clear comparison between Direct vs Regular plans increases customer retention.",
        "4. SIP Engagement: Gamify SIP milestones or alerts during market dips to maintain monthly investing discipline."
    ]
    for rec in recs:
        p_rec = tf.add_paragraph()
        p_rec.text = f"• {rec}"
        p_rec.font.size = Pt(13)
        p_rec.font.color.rgb = TEXT_LIGHT
        p_rec.space_after = Pt(8)

    # ------------------ SLIDE 12: THANK YOU ------------------
    slide12 = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide12)
    
    tb = slide12.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Questions & Answers"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "GitHub Repository: Complete ETL Pipeline & Reporting Suite\nDeliverables: Final_Report.pdf & Bluestock_MF_Presentation.pptx"
    p3.font.name = "Arial"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER
    
    prs.save(OUTPUT_PATH)
    print(f"Presentation successfully generated at {OUTPUT_PATH}")

if __name__ == "__main__":
    main()
